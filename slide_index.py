from functools import reduce
import copy

import numpy as np
import pptx
from pptx.presentation import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.fill import _GradientStop, _GradientStops
from pptx.dml.color import RGBColor


def _find_index_slide(prs: Presentation) -> int:
    for i in range(len(prs.slides)):
        if prs.slides[i].shapes.title.text == "목차":
            return i
    raise ValueError("목차가 없습니다.")


def add_index(prs: Presentation) -> list[str]:
    idx = _find_index_slide(prs)
    index_slide = prs.slides[idx]
    other_slides = [prs.slides[i] for i in range(idx + 1, len(prs.slides) - 1)]  # "목차"의 다음 슬라이드부터 "마무리" 이전 슬라이드까지

    index_list: list[tuple[str, int]] = []
    last_title = other_slides[0].shapes.title.text
    count = 1
    for slide in other_slides[1:]:
        title = slide.shapes.title.text
        if title == last_title:
            count += 1
        else:
            index_list.append((last_title, count))
            last_title = title
            count = 1
    index_list.append((title, count))
    assert reduce(lambda sum, x: sum + x[1], index_list, 0) == len(other_slides)

    index_content = index_slide.shapes.placeholders[1]
    tf = index_content.text_frame
    tf.text = index_list[0][0]

    for index in index_list[1:]:
        p = tf.add_paragraph()
        p.text = index[0]
    return index_list


def add_index_sidebar(prs: Presentation, index: list[tuple[str, int]]):
    title_and_body_layout = prs.slide_master.slide_layouts[1]
    slide_num_shape = max(title_and_body_layout.shapes, key=(lambda shape: shape.left - shape.top))  # 가장 오른쪽 위에 있는 shape
    x, y, w, h = slide_num_shape.left, slide_num_shape.top, slide_num_shape.width, slide_num_shape.height
    
    idx = _find_index_slide(prs)
    other_slides = [prs.slides[i] for i in range(1, len(prs.slides))]
    current_index = 0
    current_count = 0
    for page, slide in enumerate(other_slides):
        ys = np.linspace(y, y+h, len(index)+1)
        for i, (title, count) in enumerate(index):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, ys[i], w, ys[i+1] - ys[i])
            # 상자를 단색으로 색칠하고 그림자와 테두리를 없앤다.
            shape.fill.solid()
            shape.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_2
            shape.shadow.inherit = False
            shape.line.fill.background()
            shape.text = title
            ps = shape.text_frame.paragraphs
            ps[0].alignment = PP_ALIGN.CENTER
            ps[0].font.bold = True
            
            if page >= idx:
                if title == index[current_index][0]:
                    current_count += 1
                    if current_count > count:
                        current_index += 1
                        current_count = 0
                    else:
                        shape.fill.gradient()
                        _set_gradient(shape.fill.gradient_stops)
                        shape.text = f"\n{title}\n{current_count} / {count}"
                        ps = shape.text_frame.paragraphs
                        ps[0].font.size = Pt(10)
                        ps[1].alignment = PP_ALIGN.CENTER
                        ps[1].font.bold = True
                        ps[1].font.color.theme_color = MSO_THEME_COLOR.ACCENT_2
                        ps[2].alignment = PP_ALIGN.CENTER
                        ps[2].font.size = Pt(10)
                        ps[2].font.color.theme_color = MSO_THEME_COLOR.ACCENT_2


def _set_gradient(stops: _GradientStops):
    stop = copy.deepcopy(stops._gsLst[1])
    stops._gsLst.append(stop)
    stop = copy.deepcopy(stops._gsLst[0])
    stops._gsLst.append(stop)

    gray = RGBColor(115, 115, 115)
    stops[0].position = 0.0
    stops[0].color.rgb = gray
    stops[1].position = 0.03
    stops[1].color.theme_color = MSO_THEME_COLOR.ACCENT_4
    stops[2].position = 0.94
    stops[2].color.theme_color = MSO_THEME_COLOR.ACCENT_4
    stops[3].position = 1.0
    stops[3].color.rgb = gray
