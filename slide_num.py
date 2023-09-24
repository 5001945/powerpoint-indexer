import pptx
from pptx.slide import Slide, SlideMaster, SlideLayout
from pptx.presentation import Presentation
from pptx.shapes.base import BaseShape
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.oxml import etree
from pptx.enum.base import EnumValue
from pptx.enum.text import PP_ALIGN


def add_total_slide_num(prs: Presentation):
    slides = prs.slides
    total_num = len(slides)
    for i, slide in enumerate(slides):
        slide: Slide
        slide_num_shapes = [shape for shape in slide.shapes
                            if shape.is_placeholder and shape.placeholder_format.element.get('type') == "sldNum"]
        if slide_num_shapes:
            slide_num_shapes[-1].text = f"{i+1} / {total_num}"
            slide_num_shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


    # master: SlideMaster = prs.slide_master
    # title_and_body_layout: SlideLayout = master.slide_layouts[1]
    # slide_num_shape = max(title_and_body_layout.shapes, key=(lambda shape: shape.top + shape.left))
    # # slide_num_shape.text += f" / {total_num}"
    # slide_num_shape: LayoutPlaceholder
    # # slide_num_shape.element
    # print(slide_num_shape)
    # # print(etree.tostring(slide_num_shape.element, pretty_print=True))
    # # for shape in title_and_body_layout.shapes:
    # #     # shape: shapetree.BaseShape.
    # #     print(shape.top, shape.left)
    # #     # print(shape)
